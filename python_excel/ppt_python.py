# -- coding: utf-8 --
from pptx import Presentation

# Successfully installed Pillow-9.2.0 XlsxWriter-3.0.3 lxml-4.9.1 python-pptx-0.6.21
m_ppt = Presentation('Nginx.pptx')
print(len(m_ppt.slides))
for slide in m_ppt.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for content in paragraph.runs:
                content = content.text.replace(' ', '')
                content = content.replace('\n', '')
                content = content.replace('\t', '')
                if content:
                    print(content)